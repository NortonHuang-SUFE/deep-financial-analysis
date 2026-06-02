#!/usr/bin/env python3
"""Create a 16:9 PPTX from rendered slide images.

Use this after rendering an HTML deck into one image per slide. Images are
placed full-bleed, one per slide, preserving the visual design from HTML.
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path


def natural_key(path: Path) -> list[object]:
    return [int(part) if part.isdigit() else part.lower() for part in re.split(r"(\d+)", path.name)]


def collect_images(input_dir: Path) -> list[Path]:
    exts = {".png", ".jpg", ".jpeg"}
    return sorted((p for p in input_dir.iterdir() if p.suffix.lower() in exts), key=natural_key)


def main() -> int:
    parser = argparse.ArgumentParser(description="Build a PPTX from full-slide PNG/JPG images.")
    parser.add_argument("input_dir", help="Directory containing numbered slide images.")
    parser.add_argument("output_pptx", help="Output .pptx path.")
    parser.add_argument("--width", type=float, default=13.333333, help="Slide width in inches.")
    parser.add_argument("--height", type=float, default=7.5, help="Slide height in inches.")
    args = parser.parse_args()

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("ERROR: python-pptx is required. Install project dependencies before converting images to PPTX.", file=sys.stderr)
        return 2

    input_dir = Path(args.input_dir)
    output_pptx = Path(args.output_pptx)
    images = collect_images(input_dir)
    if not images:
        print(f"ERROR: no PNG/JPG slide images found in {input_dir}", file=sys.stderr)
        return 1

    prs = Presentation()
    prs.slide_width = Inches(args.width)
    prs.slide_height = Inches(args.height)
    blank_layout = prs.slide_layouts[6]

    for image in images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(image), 0, 0, width=prs.slide_width, height=prs.slide_height)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    print(output_pptx)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

