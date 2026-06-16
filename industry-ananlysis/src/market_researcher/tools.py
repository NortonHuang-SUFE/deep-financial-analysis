"""Local tools for the Market Researcher agent.

These tools use local Python libraries plus the installed system browser for
HTML deck rendering.
They make NO external network calls and use NO MCP servers.

Both tools are decorated with @tool so deepagents can register them
as callable agent tools alongside MCP and search tools.
"""

from __future__ import annotations

import json
import os
import html
import re
import shutil
from datetime import datetime
from pathlib import Path
from typing import Any

from langchain_core.tools import tool

from market_researcher.config import file_storage_root


# ── Helpers ───────────────────────────────────────────────────────────────────


_TASK_OUTPUT_DIRS: dict[str, Path] = {}


def _project_root() -> Path:
    return Path(__file__).resolve().parents[2]


def _workspace_root() -> Path:
    return file_storage_root()


def _resolve_output_dir(output_dir: str) -> Path:
    path = Path(output_dir)
    return path if path.is_absolute() else _workspace_root() / path


def _ensure_out_dir(output_dir: str) -> Path:
    out = _resolve_output_dir(output_dir)
    out.mkdir(parents=True, exist_ok=True)
    return out


def _timestamped_output_dir(output_dir: str) -> Path:
    """Return this process task's timestamped artifact directory under output_dir."""
    import re

    base = _resolve_output_dir(output_dir)
    if re.fullmatch(r"\d{8}-\d{6}(?:-\d+)?", base.name):
        base.mkdir(parents=True, exist_ok=True)
        return base

    key = str(base.resolve())
    existing = _TASK_OUTPUT_DIRS.get(key)
    if existing:
        existing.mkdir(parents=True, exist_ok=True)
        return existing

    timestamp = os.getenv("MARKET_RESEARCHER_OUTPUT_TIMESTAMP")
    if not timestamp:
        timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")

    candidate = base / timestamp
    if candidate.exists() and os.getenv("MARKET_RESEARCHER_OUTPUT_TIMESTAMP") is None:
        suffix = 2
        while (base / f"{timestamp}-{suffix}").exists():
            suffix += 1
        candidate = base / f"{timestamp}-{suffix}"

    candidate.mkdir(parents=True, exist_ok=True)
    _TASK_OUTPUT_DIRS[key] = candidate
    return candidate


def _slugify(text: str) -> str:
    """Very simple slug: lowercase, spaces → hyphens, drop special chars."""
    import re
    return re.sub(r"[^a-z0-9\-]", "", text.lower().replace(" ", "-"))


# ── Excel comps builder ───────────────────────────────────────────────────────


@tool
def build_comps_excel(data_json: str, sector: str, output_dir: str = "./out") -> str:
    """Build an institutional-grade comparable company analysis Excel file.

    All derived metrics (margins, multiples, statistics) are written as
    Excel formula strings so the sheet recalculates automatically when
    inputs change.  Hard-coded inputs have cell comments with the source.

    Args:
        data_json: JSON string — a list of company dicts, each with:
            {
                "ticker": "PYPL",
                "company": "PayPal Holdings",
                "revenue_ltm": 29800,           # USD millions
                "revenue_growth_pct": 8.1,       # YoY %
                "gross_profit": 14200,
                "ebitda_ltm": 5800,
                "net_income": 4200,
                "market_cap": 68000,
                "enterprise_value": 65000,
                "source": "FactSet LTM Q4 2024",  # citation for cell comment
                # optional fields:
                "fcf": 5200,
                "extra_metrics": {}              # dict of {label: value}
            }
        sector: Sector name used in the file name (e.g. "fintech-payments").
        output_dir: Base directory to write the file into. A timestamped
            subdirectory is created automatically under this directory.

    Returns:
        Relative path of the written .xlsx file.
    """
    try:
        import openpyxl
        from openpyxl.styles import (
            Alignment,
            Font,
            PatternFill,
            numbers,
        )
        from openpyxl.comments import Comment
    except ImportError as e:
        return f"ERROR: openpyxl is not installed — {e}"

    companies: list[dict[str, Any]] = json.loads(data_json)
    out_dir = _timestamped_output_dir(output_dir)
    today = datetime.now().strftime("%Y%m%d")
    slug = _slugify(sector)
    filepath = out_dir / f"comps-{slug}-{today}.xlsx"

    wb = openpyxl.Workbook()

    # ── Sheet 1: Analysis ──────────────────────────────────────────────────────
    ws = wb.active
    ws.title = "Comps"

    # Color palette
    NAVY = "1F4E79"
    LIGHT_BLUE = "D9E1F2"
    LIGHT_GRAY = "F2F2F2"
    WHITE = "FFFFFF"

    def navy_fill():
        return PatternFill("solid", fgColor=NAVY)

    def blue_fill():
        return PatternFill("solid", fgColor=LIGHT_BLUE)

    def gray_fill():
        return PatternFill("solid", fgColor=LIGHT_GRAY)

    def white_fill():
        return PatternFill("solid", fgColor=WHITE)

    def bold_white(size=12):
        return Font(name="Times New Roman", bold=True, color=WHITE, size=size)

    def bold_black(size=12):
        return Font(name="Times New Roman", bold=True, color="000000", size=size)

    def normal(size=11):
        return Font(name="Times New Roman", size=size)

    def center():
        return Alignment(horizontal="center", vertical="center")

    def left():
        return Alignment(horizontal="left", vertical="center")

    # ── Header block (rows 1-3) ────────────────────────────────────────────────
    tickers = " • ".join(
        f"{c.get('company', c.get('ticker', ''))} ({c.get('ticker', '')})"
        for c in companies
    )
    ws["A1"] = f"{sector.upper()} — COMPARABLE COMPANY ANALYSIS"
    ws["A2"] = tickers
    ws["A3"] = f"As of {datetime.now().strftime('%B %d, %Y')} | All figures in USD Millions except per-share amounts and ratios"

    for row_num in [1, 2, 3]:
        cell = ws.cell(row=row_num, column=1)
        cell.font = bold_white(12) if row_num == 1 else normal(10)
        cell.fill = navy_fill() if row_num == 1 else white_fill()
        cell.alignment = left()

    ws.merge_cells("A1:J1")
    ws.merge_cells("A2:J2")
    ws.merge_cells("A3:J3")

    # ── Section header helper ─────────────────────────────────────────────────
    def section_header(row: int, text: str, num_cols: int = 10):
        ws.cell(row=row, column=1, value=text)
        ws.cell(row=row, column=1).font = bold_white(11)
        ws.cell(row=row, column=1).fill = navy_fill()
        ws.cell(row=row, column=1).alignment = left()
        for c in range(2, num_cols + 1):
            ws.cell(row=row, column=c).fill = navy_fill()
        ws.merge_cells(
            start_row=row, start_column=1, end_row=row, end_column=num_cols
        )

    # ── Column header helper ──────────────────────────────────────────────────
    def col_header(row: int, col: int, text: str):
        cell = ws.cell(row=row, column=col, value=text)
        cell.font = bold_black(11)
        cell.fill = blue_fill()
        cell.alignment = center()

    # ── Data cell helper ──────────────────────────────────────────────────────
    def data_cell(row: int, col: int, value: Any, source: str = "", fmt: str = ""):
        cell = ws.cell(row=row, column=col, value=value)
        cell.font = normal(11)
        cell.fill = white_fill()
        cell.alignment = center()
        if fmt:
            cell.number_format = fmt
        if source:
            comment = Comment(source, "market-researcher")
            comment.width = 200
            comment.height = 60
            cell.comment = comment

    # ── Stat row helper ───────────────────────────────────────────────────────
    def stat_row(row: int, label: str, col_start: int, col_end: int, formula_fn):
        ws.cell(row=row, column=1, value=label).font = bold_black(10)
        ws.cell(row=row, column=1).fill = gray_fill()
        ws.cell(row=row, column=1).alignment = left()
        for c in range(2, col_start):
            ws.cell(row=row, column=c).fill = gray_fill()
        for c in range(col_start, col_end + 1):
            cell = ws.cell(row=row, column=c, value=formula_fn(c))
            cell.fill = gray_fill()
            cell.font = normal(10)
            cell.alignment = center()

    # ══════════════════════════════════════════════════════════════════════════
    # SECTION 1: OPERATING STATISTICS
    # Row 5: section header
    # Row 6: column headers
    # Rows 7..(7+n-1): company data
    # blank row
    # stats rows
    # ══════════════════════════════════════════════════════════════════════════
    R_OP_HDR = 5
    R_OP_COL = 6
    R_OP_DATA_START = 7
    n = len(companies)
    R_OP_DATA_END = R_OP_DATA_START + n - 1
    R_OP_BLANK = R_OP_DATA_END + 1
    R_OP_STATS_START = R_OP_BLANK + 1
    R_OP_STATS_END = R_OP_STATS_START + 4  # Max, 75th, Median, 25th, Min

    section_header(R_OP_HDR, "OPERATING STATISTICS & FINANCIAL METRICS")

    op_cols = [
        (1, "Company"),
        (2, "Revenue (LTM)"),
        (3, "Revenue Growth"),
        (4, "Gross Profit"),
        (5, "Gross Margin"),
        (6, "EBITDA (LTM)"),
        (7, "EBITDA Margin"),
        (8, "Net Income"),
        (9, "FCF"),
        (10, "FCF Margin"),
    ]
    for col, label in op_cols:
        col_header(R_OP_COL, col, label)

    for i, co in enumerate(companies):
        r = R_OP_DATA_START + i
        source = co.get("source", "")
        # Col 1: Company name
        ws.cell(row=r, column=1, value=co.get("company", co.get("ticker", ""))).font = normal(11)
        ws.cell(row=r, column=1).alignment = left()
        # Col 2: Revenue (hardcoded input)
        data_cell(r, 2, co.get("revenue_ltm", ""), source, "#,##0")
        # Col 3: Revenue Growth (hardcoded input %)
        data_cell(r, 3, co.get("revenue_growth_pct", "") / 100 if co.get("revenue_growth_pct") is not None else "", source, "0.0%")
        # Col 4: Gross Profit (hardcoded)
        data_cell(r, 4, co.get("gross_profit", ""), source, "#,##0")
        # Col 5: Gross Margin = formula D/B
        gp_col = _col_letter(4)
        rev_col = _col_letter(2)
        ws.cell(row=r, column=5, value=f"={gp_col}{r}/{rev_col}{r}").number_format = "0.0%"
        ws.cell(row=r, column=5).alignment = center()
        # Col 6: EBITDA (hardcoded)
        data_cell(r, 6, co.get("ebitda_ltm", ""), source, "#,##0")
        # Col 7: EBITDA Margin = formula F/B
        ebitda_col = _col_letter(6)
        ws.cell(row=r, column=7, value=f"={ebitda_col}{r}/{rev_col}{r}").number_format = "0.0%"
        ws.cell(row=r, column=7).alignment = center()
        # Col 8: Net Income (hardcoded)
        data_cell(r, 8, co.get("net_income", ""), source, "#,##0")
        # Col 9: FCF (hardcoded)
        data_cell(r, 9, co.get("fcf", ""), source, "#,##0")
        # Col 10: FCF Margin = formula I/B
        fcf_col = _col_letter(9)
        if co.get("fcf"):
            ws.cell(row=r, column=10, value=f"={fcf_col}{r}/{rev_col}{r}").number_format = "0.0%"
            ws.cell(row=r, column=10).alignment = center()

    # Stats rows — only for ratio/% columns (3, 5, 7, 10)
    _write_stats_block(ws, R_OP_STATS_START, R_OP_DATA_START, R_OP_DATA_END,
                       stat_cols=[3, 5, 7, 10], gray_fill=gray_fill, normal=normal)

    # ══════════════════════════════════════════════════════════════════════════
    # SECTION 2: VALUATION MULTIPLES
    # ══════════════════════════════════════════════════════════════════════════
    R_VAL_HDR = R_OP_STATS_END + 2
    R_VAL_COL = R_VAL_HDR + 1
    R_VAL_DATA_START = R_VAL_COL + 1
    R_VAL_DATA_END = R_VAL_DATA_START + n - 1
    R_VAL_BLANK = R_VAL_DATA_END + 1
    R_VAL_STATS_START = R_VAL_BLANK + 1
    R_VAL_STATS_END = R_VAL_STATS_START + 4

    section_header(R_VAL_HDR, "VALUATION MULTIPLES & INVESTMENT METRICS")

    val_cols = [
        (1, "Company"),
        (2, "Market Cap"),
        (3, "Enterprise Value"),
        (4, "EV/Revenue"),
        (5, "EV/EBITDA"),
        (6, "P/E Ratio"),
    ]
    for col, label in val_cols:
        col_header(R_VAL_COL, col, label)

    for i, co in enumerate(companies):
        r = R_VAL_DATA_START + i
        op_row = R_OP_DATA_START + i  # corresponding operating data row
        source = co.get("source", "")

        ws.cell(row=r, column=1, value=co.get("company", co.get("ticker", ""))).font = normal(11)
        ws.cell(row=r, column=1).alignment = left()

        data_cell(r, 2, co.get("market_cap", ""), source, "#,##0")
        data_cell(r, 3, co.get("enterprise_value", ""), source, "#,##0")

        ev_col = _col_letter(3)
        rev_ref = f"{_col_letter(2)}{op_row}"   # Revenue from operating section
        ebitda_ref = f"{_col_letter(6)}{op_row}"  # EBITDA from operating section
        ni_ref = f"{_col_letter(8)}{op_row}"     # Net income
        mktcap_col = _col_letter(2)

        # EV/Revenue — references operating section
        ws.cell(row=r, column=4, value=f"={ev_col}{r}/{rev_ref}").number_format = '0.0"x"'
        ws.cell(row=r, column=4).alignment = center()
        # EV/EBITDA
        ws.cell(row=r, column=5, value=f"={ev_col}{r}/{ebitda_ref}").number_format = '0.0"x"'
        ws.cell(row=r, column=5).alignment = center()
        # P/E
        if co.get("net_income") and co["net_income"] > 0:
            ws.cell(row=r, column=6, value=f"={mktcap_col}{r}/{ni_ref}").number_format = '0.0"x"'
            ws.cell(row=r, column=6).alignment = center()

    _write_stats_block(ws, R_VAL_STATS_START, R_VAL_DATA_START, R_VAL_DATA_END,
                       stat_cols=[4, 5, 6], gray_fill=gray_fill, normal=normal)

    # ══════════════════════════════════════════════════════════════════════════
    # Sheet 2: Notes & Methodology
    # ══════════════════════════════════════════════════════════════════════════
    ws_notes = wb.create_sheet("Notes")
    ws_notes["A1"] = "NOTES & METHODOLOGY"
    ws_notes["A1"].font = Font(name="Times New Roman", bold=True, size=14)
    ws_notes["A3"] = "Data Sources"
    ws_notes["A3"].font = Font(name="Times New Roman", bold=True, size=12)
    sources_used = list({c.get("source", "") for c in companies if c.get("source")})
    for idx, src in enumerate(sources_used, start=4):
        ws_notes.cell(row=idx, column=1, value=f"• {src}")
    ws_notes["A9"] = "Key Definitions"
    ws_notes["A9"].font = Font(name="Times New Roman", bold=True, size=12)
    definitions = [
        "EBITDA: Earnings Before Interest, Taxes, Depreciation & Amortization",
        "LTM: Last Twelve Months",
        "FCF: Free Cash Flow = Operating Cash Flow - Capital Expenditures",
        "Enterprise Value: Market Cap + Total Debt - Cash & Equivalents",
        "EV/Revenue: Enterprise Value / LTM Revenue",
        "EV/EBITDA: Enterprise Value / LTM EBITDA",
    ]
    for idx, d in enumerate(definitions, start=10):
        ws_notes.cell(row=idx, column=1, value=f"• {d}")

    # ── Column widths ──────────────────────────────────────────────────────────
    ws.column_dimensions["A"].width = 28
    for col_letter in ["B", "C", "D", "E", "F", "G", "H", "I", "J"]:
        ws.column_dimensions[col_letter].width = 15

    wb.save(filepath)
    return _rel(filepath)


def _col_letter(col_index: int) -> str:
    """Convert 1-based column index to Excel column letter (A, B, ..., Z, AA...)."""
    letters = ""
    while col_index > 0:
        col_index, remainder = divmod(col_index - 1, 26)
        letters = chr(65 + remainder) + letters
    return letters


def _write_stats_block(ws, stats_start_row: int, data_start: int, data_end: int,
                       stat_cols: list[int], gray_fill, normal):
    """Write Max / 75th / Median / 25th / Min rows for the given columns."""
    from openpyxl.styles import Alignment, Font, PatternFill

    stat_defs = [
        ("Maximum", "MAX"),
        ("75th Percentile", "QUARTILE({range},3)"),
        ("Median", "MEDIAN"),
        ("25th Percentile", "QUARTILE({range},1)"),
        ("Minimum", "MIN"),
    ]

    for s_idx, (label, formula_template) in enumerate(stat_defs):
        r = stats_start_row + s_idx
        cell_label = ws.cell(row=r, column=1, value=label)
        cell_label.font = Font(name="Times New Roman", bold=True, size=10)
        cell_label.fill = PatternFill("solid", fgColor="F2F2F2")
        cell_label.alignment = Alignment(horizontal="left", vertical="center")

        for c in stat_cols:
            range_str = f"{_col_letter(c)}{data_start}:{_col_letter(c)}{data_end}"
            if "QUARTILE" in formula_template:
                formula = "=" + formula_template.replace("{range}", range_str)
            else:
                formula = f"={formula_template}({range_str})"
            cell = ws.cell(row=r, column=c, value=formula)
            cell.fill = PatternFill("solid", fgColor="F2F2F2")
            cell.font = Font(name="Times New Roman", size=10)
            cell.alignment = Alignment(horizontal="center", vertical="center")
            # Copy number format from the first data cell in that column
            first_data_cell = ws.cell(row=data_start, column=c)
            cell.number_format = first_data_cell.number_format


# ── Swiss PPTX builder ────────────────────────────────────────────────────────


_SWISS_LAYOUTS = {
    "cover": "SWISS-COVER-ASCII",
    "kpi": "S02",
    "cards": "S15",
    "comparison": "S08",
    "bar_chart": "S07",
    "table": "S21",
    "closing": "SWISS-CLOSING-ASCII",
}


def _skill_dir() -> Path:
    return _project_root() / "skills" / "pptx-author"


def _rel(path: Path) -> str:
    try:
        return str(path.resolve().relative_to(_workspace_root()))
    except ValueError:
        return str(path)


def _e(value: Any) -> str:
    return html.escape("" if value is None else str(value), quote=True)


def _items(value: Any) -> list[Any]:
    if value is None:
        return []
    if isinstance(value, list):
        return value
    return [value]


def _num(value: Any) -> float:
    if isinstance(value, (int, float)):
        return float(value)
    match = re.search(r"-?\d+(?:,\d{3})*(?:\.\d+)?", str(value))
    return float(match.group(0).replace(",", "")) if match else 0.0


def _html_list(items: list[Any], color: str = "var(--text-secondary)") -> str:
    return "".join(
        f'<li style="font-size:max(18px,1.05vw);line-height:1.55;color:{color};margin-bottom:1vh">{_e(item)}</li>'
        for item in items
    )


def _chrome(idx: int, total: int, section: str) -> str:
    return (
        '<div class="chrome-min">'
        '<div class="l">MARKET RESEARCH · SWISS DECK</div>'
        f'<div class="r">{_e(section)} · {idx:02d} / {total:02d}</div>'
        "</div>"
    )


def _head(slide: dict[str, Any]) -> str:
    return (
        '<div data-anim="head" style="display:flex;flex-direction:column;gap:1.2vh;margin-bottom:3.2vh">'
        f'<div class="t-meta" style="letter-spacing:.18em;color:var(--text-helper)">{_e(slide.get("kicker", "FIELD NOTE"))}</div>'
        f'<h2 class="h-xl-zh" style="font-size:min(5.2vw,9.2vh);max-width:14ch">{_e(slide.get("title", ""))}</h2>'
        "</div>"
    )


def _render_cover(slide: dict[str, Any], idx: int, total: int, plan: dict[str, Any]) -> str:
    title = _e(slide.get("title") or plan.get("title") or "Untitled Deck").replace("\n", "<br/>")
    subtitle = slide.get("subtitle") or plan.get("subtitle", "")
    meta = slide.get("meta") or plan.get("meta", "")
    return f"""<section class="slide accent" data-layout="SWISS-COVER-ASCII" data-animate="hero">
  <div class="canvas-card">
    <canvas class="ascii-bg" aria-hidden="true"></canvas>
    {_chrome(idx, total, "COVER")}
    <div style="flex:1;padding:0;display:grid;grid-template-rows:auto 1fr auto;gap:2.6vh">
      <div data-anim="kicker" class="t-meta" style="color:rgba(255,255,255,.78);letter-spacing:.22em">{_e(slide.get("kicker", "SECTOR FIELD NOTE"))}</div>
      <h1 data-anim="title" style="font-family:var(--sans),var(--sans-zh);font-weight:200;font-size:min(10.6vw,18vh);line-height:.94;letter-spacing:-.025em;color:#fff">{title}</h1>
      <div data-anim="bottom" style="display:grid;grid-template-rows:auto auto;gap:1.5vh;border-top:1px solid rgba(255,255,255,.22);padding-top:2vh">
        <div class="lead" style="max-width:56ch;color:rgba(255,255,255,.86)">{_e(subtitle)}</div>
        <div style="display:flex;justify-content:space-between;align-items:end">
          <div class="t-meta" style="color:rgba(255,255,255,.62)">{_e(meta)}</div>
          <div class="t-meta" style="color:rgba(255,255,255,.62)">Swiss · IKB</div>
        </div>
      </div>
    </div>
  </div>
</section>"""


def _render_kpi(slide: dict[str, Any], idx: int, total: int) -> str:
    metrics = _items(slide.get("metrics"))[:6]
    cards = []
    for metric_idx, metric in enumerate(metrics):
        if isinstance(metric, dict):
            value = metric.get("value", "")
            label = metric.get("label", "")
        else:
            value = metric
            label = ""
        color = "var(--accent)" if metric_idx == 0 else "var(--text-primary)"
        cards.append(
            f'<div class="card-fill" data-anim="up" style="padding:2.2vh 1.4vw;min-height:18vh">'
            f'<div class="num-mega" style="font-size:min(4.6vw,8.2vh);color:{color}">{_e(value)}</div>'
            f'<div class="t-meta" style="margin-top:1vh;color:var(--text-secondary)">{_e(label)}</div>'
            "</div>"
        )
    bullets = _html_list(_items(slide.get("bullets")))
    return f"""<section class="slide light" data-layout="S02" data-animate="grid-reveal">
  <div class="canvas-card">
    {_chrome(idx, total, "KPI")}
    {_head(slide)}
    <div style="flex:1;padding:0;display:grid;grid-template-columns:minmax(0,1.15fr) minmax(0,.85fr);gap:3vw;align-items:start">
      <div class="grid-12" style="grid-template-columns:repeat(2,1fr);gap:16px">{''.join(cards)}</div>
      <div class="card-ink" data-anim="right" style="padding:3vh 2vw">
        <div class="t-cat" style="color:var(--accent-bright);margin-bottom:2vh">{_e(slide.get("callout", "WHY IT MATTERS"))}</div>
        <ul style="padding-left:1.2em">{bullets}</ul>
      </div>
    </div>
    <div class="t-meta" style="color:var(--text-helper);margin-top:2vh">{_e(slide.get("footnote", ""))}</div>
  </div>
</section>"""


def _render_cards(slide: dict[str, Any], idx: int, total: int) -> str:
    items = _items(slide.get("cards") or slide.get("items") or slide.get("metrics"))[:8]
    columns = min(max(len(items), 1), 4)
    cards = []
    for item_idx, item in enumerate(items):
        data = item if isinstance(item, dict) else {"title": item}
        is_accent = bool(data.get("accent")) or item_idx == 0 and slide.get("accent_first", False)
        cls = "card-accent" if is_accent else "card-fill"
        body_color = "var(--accent-on)" if is_accent else "var(--text-secondary)"
        cards.append(
            f'<div class="{cls}" data-anim="up" style="padding:2.4vh 1.5vw">'
            f'<div class="t-cat" style="margin-bottom:1.2vh">{_e(data.get("tag", data.get("label", f"{item_idx + 1:02d}")))}</div>'
            f'<h3 style="font-weight:400;font-size:max(18px,1.55vw);line-height:1.18;margin-bottom:1.2vh">{_e(data.get("title", ""))}</h3>'
            f'<p style="font-size:max(16px,.96vw);line-height:1.5;color:{body_color};font-weight:400">{_e(data.get("body", data.get("description", "")))}</p>'
            "</div>"
        )
    return f"""<section class="slide light" data-layout="S15" data-animate="grid-reveal">
  <div class="canvas-card">
    {_chrome(idx, total, "CARDS")}
    {_head(slide)}
    <div style="flex:1;padding:0;display:grid;grid-template-columns:repeat({columns},1fr);gap:16px;align-items:stretch">{''.join(cards)}</div>
    <div class="t-meta" style="color:var(--text-helper);margin-top:2vh">{_e(slide.get("footnote", ""))}</div>
  </div>
</section>"""


def _render_comparison(slide: dict[str, Any], idx: int, total: int) -> str:
    left = slide.get("left") if isinstance(slide.get("left"), dict) else {}
    right = slide.get("right") if isinstance(slide.get("right"), dict) else {}
    left_title = slide.get("left_title") or left.get("title", "Option A")
    right_title = slide.get("right_title") or right.get("title", "Option B")
    left_items = _items(slide.get("left_items") or left.get("items") or left.get("bullets"))
    right_items = _items(slide.get("right_items") or right.get("items") or right.get("bullets"))
    return f"""<section class="slide light" data-layout="S08" data-animate="split-statement">
  <div class="canvas-card">
    {_chrome(idx, total, "COMPARE")}
    {_head(slide)}
    <div style="flex:1;padding:0;display:grid;grid-template-columns:minmax(0,1fr) 1px minmax(0,1fr);gap:2.4vw;align-items:stretch">
      <div class="card-accent" data-anim="up" style="padding:4vh 2.4vw">
        <h3 class="h-md" style="font-size:min(2.8vw,5vh);margin-bottom:3vh">{_e(left_title)}</h3>
        <ul style="padding-left:1.2em">{_html_list(left_items, "var(--accent-on)")}</ul>
      </div>
      <div style="width:1px;background:var(--border-subtle)"></div>
      <div class="card-fill" data-anim="up" style="padding:4vh 2.4vw">
        <h3 class="h-md" style="font-size:min(2.8vw,5vh);margin-bottom:3vh">{_e(right_title)}</h3>
        <ul style="padding-left:1.2em">{_html_list(right_items)}</ul>
      </div>
    </div>
    <div class="t-meta" style="color:var(--text-helper);margin-top:2vh">{_e(slide.get("footnote", ""))}</div>
  </div>
</section>"""


def _render_bar_chart(slide: dict[str, Any], idx: int, total: int) -> str:
    bars = _items(slide.get("bars") or slide.get("data"))[:10]
    parsed = []
    for item in bars:
        data = item if isinstance(item, dict) else {"label": str(item), "value": 0}
        parsed.append((data.get("label", ""), _num(data.get("value", 0)), data.get("display", data.get("value", ""))))
    max_value = max([v for _, v, _ in parsed] or [1]) or 1
    rows = []
    for label, value, display in parsed:
        fill = "var(--accent)" if slide.get("highlight") and label == slide.get("highlight") else "var(--ink)"
        rows.append(
            '<div class="bar-row" data-anim="bar" style="display:grid;grid-template-columns:7em 1fr 5em;gap:1.2vw;align-items:center;margin-bottom:1.1vh">'
            f'<div class="t-cat" style="font-size:14px;color:var(--text-primary)">{_e(label)}</div>'
            f'<div style="height:2.4vh;background:var(--grey-1);position:relative"><div style="height:100%;width:{value / max_value * 100:.1f}%;background:{fill}"></div></div>'
            f'<div class="mono" style="font-size:16px;color:var(--text-secondary)">{_e(display)}</div>'
            "</div>"
        )
    callout_value = slide.get("callout_value") or (parsed[-1][2] if parsed else "")
    callout_text = slide.get("callout_text", "")
    return f"""<section class="slide light" data-layout="S07" data-animate="bar-chart">
  <div class="canvas-card">
    {_chrome(idx, total, "VALUE")}
    {_head(slide)}
    <div style="flex:1;padding:0;display:grid;grid-template-columns:minmax(0,1fr) minmax(0,.32fr);gap:3vw;align-items:start">
      <div>{''.join(rows)}</div>
      <div class="card-accent" data-anim="right" style="padding:3vh 1.8vw">
        <div class="num-mega" style="font-size:min(5vw,8.5vh)">{_e(callout_value)}</div>
        <div style="font-size:max(18px,1.05vw);line-height:1.5;margin-top:2vh;color:var(--accent-on)">{_e(callout_text)}</div>
      </div>
    </div>
    <div class="t-meta" style="color:var(--text-helper);margin-top:2vh">{_e(slide.get("footnote", ""))}</div>
  </div>
</section>"""


def _render_table(slide: dict[str, Any], idx: int, total: int) -> str:
    raw_table = slide.get("table") or slide.get("table_data") or {}
    table_data = raw_table if isinstance(raw_table, dict) else {}
    headers = _items(slide.get("headers") or table_data.get("headers"))
    rows = _items(slide.get("rows") or table_data.get("rows") or raw_table)[:8]
    header_html = "".join(f'<th style="text-align:left;padding:12px;border-bottom:2px solid var(--accent);font-size:14px">{_e(h)}</th>' for h in headers)
    rows_html = ""
    for row in rows:
        cells = row if isinstance(row, list) else _items(row)
        rows_html += "<tr>" + "".join(f'<td style="padding:12px;border-bottom:1px solid var(--border-subtle);font-size:15px;color:var(--text-secondary)">{_e(cell)}</td>' for cell in cells) + "</tr>"
    return f"""<section class="slide light" data-layout="S21" data-animate="grid-reveal">
  <div class="canvas-card">
    {_chrome(idx, total, "TABLE")}
    {_head(slide)}
    <div style="flex:1;padding:0;display:grid;align-items:start">
      <table data-anim="up" style="width:100%;border-collapse:collapse;font-family:var(--sans),var(--sans-zh);background:var(--paper)">
        <thead><tr>{header_html}</tr></thead>
        <tbody>{rows_html}</tbody>
      </table>
    </div>
    <div class="t-meta" style="color:var(--text-helper);margin-top:2vh">{_e(slide.get("footnote", ""))}</div>
  </div>
</section>"""


def _render_closing(slide: dict[str, Any], idx: int, total: int, plan: dict[str, Any]) -> str:
    title = _e(slide.get("title", "Closing")).replace("\n", "<br/>")
    takeaways = _items(slide.get("takeaways") or slide.get("items"))[:3]
    takeaway_html = []
    for takeaway_idx, item in enumerate(takeaways):
        data = item if isinstance(item, dict) else {"title": item}
        number = data.get("number", f"{takeaway_idx + 1:02d}")
        color = "var(--accent)" if takeaway_idx == len(takeaways) - 1 else "var(--text-primary)"
        takeaway_html.append(
            '<div style="display:grid;grid-template-columns:auto 1fr;gap:2vw;align-items:start;padding:2.4vh 0;border-top:1px solid var(--border-subtle)">'
            f'<div style="font-family:var(--sans);font-weight:200;font-size:min(4vw,7vh);line-height:.9;color:{color}">{_e(number)}</div>'
            f'<div><h3 style="font-weight:400;font-size:max(18px,1.6vw);line-height:1.2;color:{color};margin-bottom:.8vh">{_e(data.get("title", ""))}</h3>'
            f'<p style="font-size:max(16px,.95vw);line-height:1.55;color:var(--text-secondary);font-weight:400">{_e(data.get("body", data.get("description", "")))}</p></div>'
            "</div>"
        )
    return f"""<section class="slide split" data-layout="SWISS-CLOSING-ASCII" data-animate="split-statement">
  <div class="canvas-card">
    <div class="split-half">
      <div class="half b-accent" style="padding:5.6vh 3.6vw 4.4vh;justify-content:space-between;position:relative;overflow:hidden">
        <canvas class="ascii-bg" aria-hidden="true"></canvas>
        {_chrome(idx, total, "CLOSING")}
        <div data-anim="manifesto">
          <div class="t-meta" style="color:rgba(255,255,255,.78);letter-spacing:.22em;margin-bottom:1.8vh">{_e(slide.get("kicker", "CONCLUSION"))}</div>
          <h2 style="font-size:min(7.4vw,13vh);line-height:.94;letter-spacing:-.025em;font-weight:200;color:#fff">{title}</h2>
          <div style="font-size:max(18px,1.1vw);line-height:1.55;color:rgba(255,255,255,.86);font-weight:400;max-width:36ch;margin-top:2vh">{_e(slide.get("subtitle", plan.get("subtitle", "")))}</div>
        </div>
        <div class="t-meta" style="color:rgba(255,255,255,.62);border-top:1px solid rgba(255,255,255,.22);padding-top:2vh">{_e(slide.get("footnote", ""))}</div>
      </div>
      <div class="half" style="padding:5.6vh 3.6vw 4.4vh;justify-content:space-between">
        <div class="chrome-min"><div class="l">TAKEAWAYS</div><div class="r">03 RULES</div></div>
        <div data-anim="rules" style="display:flex;flex-direction:column;gap:0">{''.join(takeaway_html)}</div>
        <div class="t-meta" style="color:var(--text-helper);text-align:right">END OF DECK</div>
      </div>
    </div>
  </div>
</section>"""


def _render_swiss_slide(slide: dict[str, Any], idx: int, total: int, plan: dict[str, Any]) -> str:
    layout = slide.get("layout", "")
    if layout not in _SWISS_LAYOUTS:
        raise ValueError(f"Unsupported Swiss PPT layout: {layout}. Use one of: {', '.join(_SWISS_LAYOUTS)}")
    if layout == "cover":
        return _render_cover(slide, idx, total, plan)
    if layout == "kpi":
        return _render_kpi(slide, idx, total)
    if layout == "cards":
        return _render_cards(slide, idx, total)
    if layout == "comparison":
        return _render_comparison(slide, idx, total)
    if layout == "bar_chart":
        return _render_bar_chart(slide, idx, total)
    if layout == "table":
        return _render_table(slide, idx, total)
    return _render_closing(slide, idx, total, plan)


def _build_swiss_html(plan: dict[str, Any]) -> str:
    skill_dir = _skill_dir()
    template_path = skill_dir / "assets" / "template-swiss.html"
    if not template_path.exists():
        raise FileNotFoundError(f"Swiss template not found: {template_path}")

    slides = list(plan.get("slides") or [])
    if not slides:
        raise ValueError("slides_json must contain at least one slide.")
    if slides[0].get("layout") != "cover":
        slides.insert(0, {
            "layout": "cover",
            "title": plan.get("title", "Untitled Deck"),
            "subtitle": plan.get("subtitle", ""),
            "meta": plan.get("meta", ""),
        })

    rendered = "\n\n".join(_render_swiss_slide(slide, idx + 1, len(slides), plan) for idx, slide in enumerate(slides))
    template = template_path.read_text(encoding="utf-8")
    template = re.sub(
        r'^\s*<link[^>]+https://fonts\.(?:googleapis|gstatic)\.com[^>]*>\n?',
        "",
        template,
        flags=re.MULTILINE,
    )
    template = re.sub(
        r'^\s*<script src="https://unpkg\.com/lucide[^"]+"></script>\n?\s*<script>lucide\.createIcons\(\);</script>\n?',
        "<script>window.lucide={createIcons(){}};</script>\n",
        template,
        flags=re.MULTILINE,
    )
    template = template.replace(
        "motion = await import('https://cdn.jsdelivr.net/npm/motion@11.11.17/+esm');",
        "throw e1;",
    )
    template = template.replace(
        "[motion] local + CDN both failed, disabling animations",
        "[motion] local motion failed, disabling animations",
    )
    title = _e(plan.get("title", "Swiss PPT"))
    template = template.replace("[必填] 替换为 PPT 标题 · Deck Title", f"{title} · Swiss Deck")
    start_marker = "<!-- ============================================================\n     SLIDES 插入区"
    end_marker = '\n</div>\n\n<div id="nav"></div>'
    start = template.index(start_marker)
    end = template.index(end_marker, start)
    return template[:start] + rendered + template[end:]


def _find_browser() -> str | None:
    for candidate in [
        os.getenv("PPTX_RENDER_BROWSER"),
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
        "/Applications/Chromium.app/Contents/MacOS/Chromium",
        "/Applications/Brave Browser.app/Contents/MacOS/Brave Browser",
    ]:
        if candidate and Path(candidate).exists():
            return candidate
    return None


def _render_html_to_images(html_path: Path, image_dir: Path) -> int:
    try:
        from playwright.sync_api import sync_playwright
    except ImportError as e:
        raise ImportError("playwright is required for Swiss PPT rendering. Install project dependencies.") from e

    image_dir.mkdir(parents=True, exist_ok=True)
    for old_image in image_dir.glob("*.png"):
        old_image.unlink()

    with sync_playwright() as p:
        launch_kwargs: dict[str, Any] = {"headless": True}
        browser_path = _find_browser()
        if browser_path:
            launch_kwargs["executable_path"] = browser_path
        browser = p.chromium.launch(**launch_kwargs)
        page = browser.new_page(device_scale_factor=1, viewport={"width": 1920, "height": 1080})
        page.route("http://*/*", lambda route: route.abort())
        page.route("https://*/*", lambda route: route.abort())
        page.add_init_script("localStorage.setItem('guizang-ppt-low-power','1');")
        page.goto(html_path.resolve().as_uri(), wait_until="domcontentloaded", timeout=60000)
        page.wait_for_timeout(400)
        page.evaluate(
            """() => {
              document.body.classList.add('low-power');
              const deck = document.querySelector('#deck');
              if (deck) deck.style.transition = 'none';
              document.querySelectorAll('[data-anim]').forEach(el => {
                el.style.opacity = '1';
                el.style.transform = 'none';
              });
            }"""
        )
        slide_count = page.locator(".slide").count()
        if slide_count == 0:
            browser.close()
            raise ValueError("No .slide elements found in generated HTML.")
        for idx in range(slide_count):
            page.evaluate(
                """idx => {
                  const deck = document.querySelector('#deck');
                  if (deck) deck.style.transform = `translateX(${-idx * 100}vw)`;
                  window.dispatchEvent(new Event('resize'));
                }""",
                idx,
            )
            page.wait_for_timeout(150)
            page.screenshot(path=str(image_dir / f"{idx + 1:02d}.png"), full_page=False)
        browser.close()
    return slide_count


def _images_to_pptx(image_dir: Path, pptx_path: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as e:
        raise ImportError("python-pptx is required for Swiss PPT packaging.") from e

    images = sorted(image_dir.glob("*.png"))
    if not images:
        raise ValueError(f"No rendered slide images found in {image_dir}")

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    for image in images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(image), 0, 0, width=prs.slide_width, height=prs.slide_height)
    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(pptx_path)


@tool
def build_pptx(slides_json: str, output_name: str, output_dir: str = "./out",
               template_path: str = "./templates/firm-template.pptx") -> str:
    """Build a Swiss International Style PPTX from a structured deck plan.

    The public name remains `build_pptx`, but this is the only PPT pipeline:
    it writes a Swiss HTML source deck, renders the HTML to slide PNGs with
    Playwright + system Chrome/Edge, and packages those PNGs into the final
    PPTX. `template_path` is accepted for backwards call compatibility and is
    intentionally ignored.

    Args:
        slides_json: JSON string with:
            {
                "title": "Deck title",
                "subtitle": "Optional deck subtitle",
                "meta": "Optional date/source",
                "slides": [
                    {
                        "layout": "cover|kpi|cards|comparison|bar_chart|table|closing",
                        "title": "Insight headline",
                        "kicker": "SECTION · 01",
                        "...": "layout-specific fields"
                    }
                ]
            }
        output_name: Filename without extension.
        output_dir: Base directory or existing timestamp directory. The tool
            writes artifacts under `<timestamp>/ppt/`.
        template_path: Deprecated; ignored.

    Returns:
        Human-readable artifact summary with the PPTX path first.
    """
    del template_path
    try:
        plan: dict[str, Any] = json.loads(slides_json)
        out_dir = _timestamped_output_dir(output_dir)
        ppt_dir = out_dir / "ppt"
        assets_dir = ppt_dir / "assets"
        rendered_dir = ppt_dir / "rendered-slides"
        ppt_dir.mkdir(parents=True, exist_ok=True)
        assets_dir.mkdir(parents=True, exist_ok=True)

        motion_src = _skill_dir() / "assets" / "motion.min.js"
        if motion_src.exists():
            shutil.copy2(motion_src, assets_dir / "motion.min.js")

        html_path = ppt_dir / "index.html"
        html_path.write_text(_build_swiss_html(plan), encoding="utf-8")

        rendered_count = _render_html_to_images(html_path, rendered_dir)
        safe_output_name = output_name[:-5] if output_name.lower().endswith(".pptx") else output_name
        pptx_path = ppt_dir / f"{safe_output_name}.pptx"
        _images_to_pptx(rendered_dir, pptx_path)

        manifest = {
            "pptx": _rel(pptx_path),
            "html": _rel(html_path),
            "rendered_slides_dir": _rel(rendered_dir),
            "slide_count": rendered_count,
            "layouts": [slide.get("layout") for slide in plan.get("slides", [])],
            "mode": "swiss-html-rendered-to-pptx",
        }
        manifest_path = ppt_dir / "manifest.json"
        manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

        return (
            f"PPTX: {_rel(pptx_path)}\n"
            f"HTML: {_rel(html_path)}\n"
            f"Rendered slides: {rendered_count} ({_rel(rendered_dir)})\n"
            f"Manifest: {_rel(manifest_path)}"
        )
    except Exception as e:
        return f"ERROR: build_pptx failed — {type(e).__name__}: {e}"
